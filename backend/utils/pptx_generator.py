import os
import tempfile
from dataclasses import dataclass
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass
class ThemeSpec:
    name: str
    background: RGBColor
    surface: RGBColor
    surface_alt: RGBColor
    accent: RGBColor
    accent_alt: RGBColor
    text: RGBColor
    muted: RGBColor
    inverse_text: RGBColor
    title_font: str = "Aptos Display"
    body_font: str = "Aptos"


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.replace("#", "")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def resolve_theme(theme_vibe: str) -> ThemeSpec:
    vibe = (theme_vibe or "").lower()

    if "google" in vibe or "startup" in vibe or "visionary" in vibe or "bold" in vibe:
        return ThemeSpec(
            name="Catalyst",
            background=rgb("#F7F9FC"),
            surface=rgb("#FFFFFF"),
            surface_alt=rgb("#E8F0FE"),
            accent=rgb("#1A73E8"),
            accent_alt=rgb("#34A853"),
            text=rgb("#0F172A"),
            muted=rgb("#64748B"),
            inverse_text=rgb("#FFFFFF"),
        )

    if "apple" in vibe or "minimal" in vibe or "clean" in vibe or "story" in vibe:
        return ThemeSpec(
            name="Monochrome",
            background=rgb("#F2F2F0"),
            surface=rgb("#FFFFFF"),
            surface_alt=rgb("#E5E7EB"),
            accent=rgb("#111111"),
            accent_alt=rgb("#6B7280"),
            text=rgb("#111111"),
            muted=rgb("#6B7280"),
            inverse_text=rgb("#FFFFFF"),
        )

    if "cyber" in vibe or "hacker" in vibe or "matrix" in vibe or "technical" in vibe or "deep" in vibe:
        return ThemeSpec(
            name="Signal Grid",
            background=rgb("#08111F"),
            surface=rgb("#101A30"),
            surface_alt=rgb("#16233E"),
            accent=rgb("#12E7F2"),
            accent_alt=rgb("#FF4ECD"),
            text=rgb("#EAFBFF"),
            muted=rgb("#93A4BF"),
            inverse_text=rgb("#08111F"),
        )

    return ThemeSpec(
        name="Executive Blueprint",
        background=rgb("#0F172A"),
        surface=rgb("#111C33"),
        surface_alt=rgb("#1E293B"),
        accent=rgb("#38BDF8"),
        accent_alt=rgb("#F59E0B"),
        text=rgb("#F8FAFC"),
        muted=rgb("#94A3B8"),
        inverse_text=rgb("#0F172A"),
    )


def accent_for_tone(theme: ThemeSpec, tone: Optional[str] = None) -> RGBColor:
    return theme.accent_alt if tone == "secondary" else theme.accent


def style_paragraph(paragraph, theme: ThemeSpec, font_size: int, color: RGBColor = None, bold: bool = False, font_name: str = None, align=PP_ALIGN.LEFT):
    paragraph.alignment = align
    if not paragraph.runs:
        paragraph.text = paragraph.text or " "

    for run in paragraph.runs:
        run.font.name = font_name or theme.body_font
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color or theme.text


def ensure_min_length(length, minimum=Inches(0.08)):
    return length if length > minimum else minimum


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    theme: ThemeSpec,
    font_size: int,
    color: RGBColor = None,
    bold: bool = False,
    font_name: str = None,
    align=PP_ALIGN.LEFT,
):
    width = ensure_min_length(width)
    height = ensure_min_length(height)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.clear()

    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    style_paragraph(paragraph, theme, font_size, color=color, bold=bold, font_name=font_name, align=align)
    return text_box


def add_body_bullets(slide, left, top, width, height, bullets: List[str], theme: ThemeSpec, font_size: int = 14):
    if not bullets:
        return

    width = ensure_min_length(width)
    height = ensure_min_length(height)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        style_paragraph(paragraph, theme, font_size, color=theme.text)
        paragraph.text = f"• {paragraph.text}"


def add_card(slide, left, top, width, height, title: str, body: str, theme: ThemeSpec, accent: RGBColor = None, compact: bool = False):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = theme.surface
    card.fill.transparency = 0.05 if theme.background != theme.surface else 0
    card.line.color.rgb = accent or theme.accent
    card.line.transparency = 0.45

    tf = card.text_frame
    tf.clear()
    inset = Inches(0.16 if compact else 0.22)
    tf.margin_left = inset
    tf.margin_top = inset
    tf.margin_right = inset
    tf.margin_bottom = inset
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True

    p_title = tf.paragraphs[0]
    p_title.text = title
    style_paragraph(p_title, theme, 10 if compact else 14, color=accent or theme.accent, bold=True)
    p_title.space_after = Pt(4 if compact else 8)

    p_body = tf.add_paragraph()
    p_body.text = body
    style_paragraph(p_body, theme, 11 if compact else 14, color=theme.text, font_name=theme.body_font)

    return card


def add_metric_card(slide, left, top, width, height, metric: Dict[str, str], theme: ThemeSpec, accent: RGBColor = None, compact: bool = False):
    accent_color = accent or theme.accent
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = theme.surface
    card.line.color.rgb = accent_color
    card.line.transparency = 0.35

    tf = card.text_frame
    tf.clear()
    inset = Inches(0.16 if compact else 0.22)
    tf.margin_left = inset
    tf.margin_top = inset
    tf.margin_right = inset
    tf.margin_bottom = inset
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True

    p_label = tf.paragraphs[0]
    p_label.text = metric.get("label", "Metric")
    style_paragraph(p_label, theme, 10 if compact else 12, color=theme.muted, bold=True)
    p_label.space_after = Pt(2 if compact else 6)

    p_value = tf.add_paragraph()
    p_value.text = metric.get("value", "Signal")
    style_paragraph(p_value, theme, 18 if compact else 30, color=accent_color, bold=True, font_name=theme.title_font)
    p_value.space_after = Pt(2 if compact else 6)

    detail = metric.get("detail", "")
    if detail:
        p_detail = tf.add_paragraph()
        p_detail.text = detail
        style_paragraph(p_detail, theme, 10 if compact else 13, color=theme.text)

    return card


def apply_background(slide, width, height, theme: ThemeSpec):
    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, width, height)
    background.fill.solid()
    background.fill.fore_color.rgb = theme.background
    background.line.fill.background()

    orb_one = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.15), Inches(0.12), Inches(2.75), Inches(2.75))
    orb_one.fill.solid()
    orb_one.fill.fore_color.rgb = theme.accent
    orb_one.fill.transparency = 0.84
    orb_one.line.fill.background()

    orb_two = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, width - Inches(2.65), height - Inches(2.35), Inches(2.25), Inches(2.25))
    orb_two.fill.solid()
    orb_two.fill.fore_color.rgb = theme.accent_alt
    orb_two.fill.transparency = 0.9
    orb_two.line.fill.background()


def add_slide_chrome(slide, slide_data: Dict[str, Any], theme: ThemeSpec, width, height, slide_number: int):
    add_text_box(slide, Inches(0.7), Inches(0.45), Inches(2.6), Inches(0.28), slide_data.get("section_label", "Strategy").upper(), theme, 10, color=theme.accent, bold=True)
    add_text_box(slide, Inches(0.7), Inches(0.78), Inches(10), Inches(0.48), slide_data.get("title", "Executive Insight"), theme, 13, color=theme.muted, bold=True)
    add_text_box(slide, Inches(0.7), Inches(1.18), Inches(8.6), Inches(0.74), slide_data.get("headline", slide_data.get("title", "")), theme, 26, color=theme.text, bold=True, font_name=theme.title_font)
    add_text_box(slide, Inches(0.7), Inches(1.95), Inches(8.8), Inches(0.56), slide_data.get("subheadline", ""), theme, 14, color=theme.muted)
    accent_rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.42), Inches(1.15), Inches(0.06))
    accent_rule.fill.solid()
    accent_rule.fill.fore_color.rgb = theme.accent
    accent_rule.line.fill.background()
    footer = add_text_box(slide, width - Inches(1.2), height - Inches(0.45), Inches(0.6), Inches(0.18), f"{slide_number}", theme, 10, color=theme.muted, align=PP_ALIGN.RIGHT)
    footer.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM


def normalize_cards(slide_data: Dict[str, Any]) -> List[Dict[str, str]]:
    cards = slide_data.get("cards") or []
    if cards:
        return cards
    return [
        {"title": f"Focus {index + 1}", "body": bullet}
        for index, bullet in enumerate((slide_data.get("bullets") or [])[:4])
    ]


def normalize_metrics(slide_data: Dict[str, Any]) -> List[Dict[str, str]]:
    metrics = slide_data.get("metrics") or []
    if metrics:
        return metrics
    bullets = slide_data.get("bullets") or []
    return [
        {"label": f"Signal {index + 1}", "value": bullet[:36], "detail": slide_data.get("subheadline", "")}
        for index, bullet in enumerate(bullets[:3])
    ]


def normalize_steps(slide_data: Dict[str, Any]) -> List[str]:
    steps = slide_data.get("flow_steps") or []
    if steps:
        return steps
    return (slide_data.get("bullets") or [])[:4]


def get_render_payload(slide_data: Dict[str, Any]) -> Dict[str, Any]:
    return slide_data.get("render_payload") or {}


def get_feature_cards(slide_data: Dict[str, Any]) -> List[Dict[str, str]]:
    payload_cards = get_render_payload(slide_data).get("feature_cards") or []
    if payload_cards:
        return payload_cards[:4]
    return normalize_cards(slide_data)[:4]


def get_metric_cards(slide_data: Dict[str, Any]) -> List[Dict[str, str]]:
    payload_metrics = get_render_payload(slide_data).get("metric_cards") or []
    if payload_metrics:
        return payload_metrics[:3]
    return normalize_metrics(slide_data)[:3]


def get_step_cards(slide_data: Dict[str, Any], label_builder) -> List[Dict[str, str]]:
    payload_steps = get_render_payload(slide_data).get("step_cards") or []
    if payload_steps:
        return payload_steps
    return [
        {"title": label_builder(index), "body": step, "tone": "primary" if index % 2 == 0 else "secondary"}
        for index, step in enumerate(normalize_steps(slide_data))
    ]


def get_comparison_cards(slide_data: Dict[str, Any]) -> List[Dict[str, str]]:
    payload_cards = get_render_payload(slide_data).get("comparison_cards") or []
    if len(payload_cards) >= 2:
        return payload_cards[:2]

    cards = normalize_cards(slide_data)[:2]
    left_card = cards[0] if cards else {"title": "Current State", "body": (slide_data.get("bullets") or ["Current position"])[0], "tone": "primary"}
    right_card = cards[1] if len(cards) > 1 else {"title": "Future State", "body": (slide_data.get("bullets") or ["Future position"])[-1], "tone": "secondary"}
    return [left_card, right_card]


def get_supporting_card(slide_data: Dict[str, Any], fallback_title: str, fallback_body: str) -> Dict[str, str]:
    payload_card = get_render_payload(slide_data).get("supporting_card")
    if payload_card:
        return payload_card
    return {"title": fallback_title, "body": fallback_body, "tone": "secondary"}


def get_supporting_cards(slide_data: Dict[str, Any], label_builder) -> List[Dict[str, str]]:
    payload_cards = get_render_payload(slide_data).get("supporting_cards") or []
    if payload_cards:
        return payload_cards[:3]
    return [
        {"title": label_builder(index), "body": bullet, "tone": "primary" if index % 2 == 0 else "secondary"}
        for index, bullet in enumerate((slide_data.get("bullets") or [])[:3])
    ]


def render_hero(slide, slide_data: Dict[str, Any], theme: ThemeSpec):
    left_width = Inches(6.1)
    right_x = Inches(7.15)
    payload = get_render_payload(slide_data)
    hero_quote = payload.get("lead_quote") or slide_data.get("quote") or slide_data.get("headline", "")
    hero_bullets = payload.get("bullet_points") or (slide_data.get("bullets") or [])[:4]
    thesis_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.75), left_width, Inches(3.15))
    thesis_panel.fill.solid()
    thesis_panel.fill.fore_color.rgb = theme.surface
    thesis_panel.fill.transparency = 0.03
    thesis_panel.line.color.rgb = theme.accent
    thesis_panel.line.transparency = 0.55

    add_text_box(slide, Inches(0.98), Inches(3.02), Inches(5.45), Inches(1.0), hero_quote, theme, 25, color=theme.text, bold=True, font_name=theme.title_font)
    add_body_bullets(slide, Inches(0.98), Inches(4.18), Inches(5.15), Inches(1.28), hero_bullets[:3], theme, font_size=13)

    cards = get_feature_cards(slide_data)[:2]
    metrics = get_metric_cards(slide_data)[:2]

    hero_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, right_x, Inches(2.6), Inches(2.95), Inches(3.35))
    hero_panel.fill.solid()
    hero_panel.fill.fore_color.rgb = theme.surface_alt
    hero_panel.line.color.rgb = theme.accent_alt
    hero_panel.line.transparency = 0.45

    add_text_box(slide, right_x + Inches(0.22), Inches(2.82), Inches(2.5), Inches(0.3), slide_data.get("accent", "Strategic Edge").upper(), theme, 10, color=theme.accent_alt, bold=True)

    if metrics:
        for index, metric in enumerate(metrics):
            add_metric_card(
                slide,
                right_x + Inches(0.16),
                Inches(3.18 + (index * 1.18)),
                Inches(2.6),
                Inches(1.1),
                metric,
                theme,
                accent=accent_for_tone(theme, metric.get("tone")),
                compact=True,
            )
    else:
        for index, card in enumerate(cards):
            add_card(
                slide,
                right_x + Inches(0.16),
                Inches(3.22 + (index * 1.18)),
                Inches(2.6),
                Inches(1.04),
                card.get("title", "Insight"),
                card.get("body", ""),
                theme,
                accent=accent_for_tone(theme, card.get("tone")),
                compact=True,
            )


def render_insight_grid(slide, slide_data: Dict[str, Any], theme: ThemeSpec):
    cards = get_feature_cards(slide_data)[:4]
    positions = [
        (Inches(0.7), Inches(2.7)),
        (Inches(5.12), Inches(2.7)),
        (Inches(0.7), Inches(4.78)),
        (Inches(5.12), Inches(4.78)),
    ]
    for index, card_data in enumerate(cards):
        left, top = positions[index]
        add_card(
            slide,
            left,
            top,
            Inches(3.95),
            Inches(1.72),
            card_data.get("title", "Insight"),
            card_data.get("body", ""),
            theme,
            accent=accent_for_tone(theme, card_data.get("tone")),
        )


def render_metrics_band(slide, slide_data: Dict[str, Any], theme: ThemeSpec):
    metrics = get_metric_cards(slide_data)[:3]
    card_width = Inches(3.02)
    for index, metric in enumerate(metrics):
        add_metric_card(
            slide,
            Inches(0.7) + (card_width + Inches(0.2)) * index,
            Inches(2.75),
            card_width,
            Inches(1.82),
            metric,
            theme,
            accent=accent_for_tone(theme, metric.get("tone")),
        )

    supporting_card = get_supporting_card(
        slide_data,
        slide_data.get("quote") or "Why this matters now",
        " • ".join((slide_data.get("bullets") or [])[:3]) or slide_data.get("subheadline", ""),
    )

    add_card(
        slide,
        Inches(0.7),
        Inches(4.95),
        Inches(9.7),
        Inches(1.55),
        supporting_card.get("title", "Why this matters now"),
        supporting_card.get("body", ""),
        theme,
        accent=accent_for_tone(theme, supporting_card.get("tone")),
    )


def render_process_flow(slide, slide_data: Dict[str, Any], theme: ThemeSpec):
    steps = get_step_cards(slide_data, lambda index: str(index + 1))[:4]
    start_left = Inches(0.85)
    step_width = Inches(2.15)
    step_height = Inches(1.25)
    gap = Inches(0.25)
    line_y = Inches(4.02)

    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.6), line_y, Inches(8.9), line_y)
    connector.line.color.rgb = theme.accent
    connector.line.transparency = 0.4
    connector.line.width = Pt(2)

    for index, step in enumerate(steps):
        left = start_left + index * (step_width + gap)
        accent = accent_for_tone(theme, step.get("tone"))
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(3.1), step_width, step_height)
        card.fill.solid()
        card.fill.fore_color.rgb = theme.surface
        card.line.color.rgb = accent
        card.line.transparency = 0.35

        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.16), Inches(3.26), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        
        btf = badge.text_frame
        btf.clear()
        btf.margin_left = 0
        btf.margin_right = 0
        btf.margin_top = 0
        btf.margin_bottom = 0
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.text = step.get("title", str(index + 1))
        style_paragraph(bp, theme, 10, color=theme.inverse_text, bold=True, align=PP_ALIGN.CENTER)

        tf = card.text_frame
        tf.clear()
        tf.margin_left = Inches(0.16)
        tf.margin_top = Inches(0.65)
        tf.margin_right = Inches(0.16)
        tf.margin_bottom = Inches(0.16)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = step.get("body", "")
        style_paragraph(p, theme, 11, color=theme.text)

    supporting_card = get_supporting_card(
        slide_data,
        slide_data.get("accent") or "Execution Signal",
        " • ".join((slide_data.get("bullets") or [])[:3]) or slide_data.get("subheadline", ""),
    )

    add_card(
        slide,
        Inches(0.7),
        Inches(5.18),
        Inches(9.7),
        Inches(1.28),
        supporting_card.get("title", slide_data.get("accent") or "Execution Signal"),
        supporting_card.get("body", ""),
        theme,
        accent=accent_for_tone(theme, supporting_card.get("tone")),
        compact=True,
    )


def render_comparison(slide, slide_data: Dict[str, Any], theme: ThemeSpec):
    left_card, right_card = get_comparison_cards(slide_data)
    supporting_cards = get_supporting_cards(slide_data, lambda index: f"Move {index + 1}")

    add_card(slide, Inches(0.7), Inches(2.75), Inches(4.6), Inches(2.55), left_card.get("title", "Current State"), left_card.get("body", ""), theme, accent=accent_for_tone(theme, left_card.get("tone")))
    add_card(slide, Inches(5.8), Inches(2.75), Inches(4.6), Inches(2.55), right_card.get("title", "Future State"), right_card.get("body", ""), theme, accent=accent_for_tone(theme, right_card.get("tone")))

    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(4.85), Inches(3.62), Inches(0.8), Inches(0.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme.accent_alt
    accent_bar.line.fill.background()

    for index, card in enumerate(supporting_cards[:3]):
        add_card(
            slide,
            Inches(0.7) + Inches(3.32) * index,
            Inches(5.42),
            Inches(3.08),
            Inches(1.06),
            card.get("title", f"Move {index + 1}"),
            card.get("body", ""),
            theme,
            accent=accent_for_tone(theme, card.get("tone")),
            compact=True,
        )


def render_roadmap(slide, slide_data: Dict[str, Any], theme: ThemeSpec):
    milestones = get_step_cards(slide_data, lambda index: f"Phase {index + 1}")[:3]
    supporting_cards = get_supporting_cards(slide_data, lambda _index: "Checkpoint")
    left_positions = [Inches(0.8), Inches(3.85), Inches(6.9)]

    track = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.4), Inches(4.55), Inches(9.4), Inches(4.55))
    track.line.color.rgb = theme.accent
    track.line.transparency = 0.45
    track.line.width = Pt(2)

    for index, milestone in enumerate(milestones):
        center = left_positions[index] + Inches(1.1)
        accent = accent_for_tone(theme, milestone.get("tone"))
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, center, Inches(4.32), Inches(0.45), Inches(0.45))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        add_card(slide, left_positions[index], Inches(2.95), Inches(2.45), Inches(1.3), milestone.get("title", f"Phase {index + 1}"), milestone.get("body", ""), theme, accent=accent, compact=True)

    for index, card in enumerate(supporting_cards[:3]):
        add_card(
            slide,
            Inches(0.7) + Inches(3.32) * index,
            Inches(5.3),
            Inches(3.08),
            Inches(1.06),
            card.get("title", "Checkpoint"),
            card.get("body", ""),
            theme,
            accent=accent_for_tone(theme, card.get("tone")),
            compact=True,
        )


def render_closing(slide, slide_data: Dict[str, Any], theme: ThemeSpec):
    payload = get_render_payload(slide_data)
    priorities = get_supporting_cards(slide_data, lambda index: f"Priority {index + 1}")
    add_text_box(slide, Inches(0.9), Inches(2.65), Inches(8.8), Inches(0.92), payload.get("lead_quote") or slide_data.get("quote") or slide_data.get("headline", ""), theme, 34, color=theme.text, bold=True, font_name=theme.title_font)
    add_text_box(slide, Inches(0.92), Inches(3.7), Inches(8.6), Inches(0.6), slide_data.get("subheadline", ""), theme, 16, color=theme.muted)

    for index, card in enumerate(priorities[:3]):
        add_card(
            slide,
            Inches(0.9) + Inches(3.15) * index,
            Inches(4.72),
            Inches(2.85),
            Inches(1.2),
            card.get("title", f"Priority {index + 1}"),
            card.get("body", ""),
            theme,
            accent=accent_for_tone(theme, card.get("tone")),
            compact=True,
        )


def render_content_slide(slide, slide_data: Dict[str, Any], theme: ThemeSpec, width, height, slide_number: int):
    apply_background(slide, width, height, theme)
    add_slide_chrome(slide, slide_data, theme, width, height, slide_number)

    layout = slide_data.get("layout_style", "insight-grid")
    if layout == "hero":
        render_hero(slide, slide_data, theme)
    elif layout == "metrics-band":
        render_metrics_band(slide, slide_data, theme)
    elif layout == "process-flow":
        render_process_flow(slide, slide_data, theme)
    elif layout == "comparison":
        render_comparison(slide, slide_data, theme)
    elif layout == "roadmap":
        render_roadmap(slide, slide_data, theme)
    elif layout == "closing":
        render_closing(slide, slide_data, theme)
    else:
        render_insight_grid(slide, slide_data, theme)


def render_cover_slide(prs: Presentation, presentation_json: Dict[str, Any], theme: ThemeSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    height = prs.slide_height
    apply_background(slide, width, height, theme)

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.85), Inches(5.3), Inches(5.7))
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme.surface
    panel.line.color.rgb = theme.accent
    panel.line.transparency = 0.6

    cover_payload = presentation_json.get("cover_payload") or {}
    add_text_box(slide, Inches(1.05), Inches(1.2), Inches(4.3), Inches(0.3), (cover_payload.get("eyebrow") or "EXECUTIVE NARRATIVE DECK").upper(), theme, 11, color=theme.accent, bold=True)
    add_text_box(slide, Inches(1.05), Inches(1.7), Inches(4.45), Inches(1.4), presentation_json.get("deck_title", "Executive Strategy Deck"), theme, 32, color=theme.text, bold=True, font_name=theme.title_font)
    add_text_box(slide, Inches(1.05), Inches(3.25), Inches(4.3), Inches(0.8), presentation_json.get("deck_subtitle", ""), theme, 16, color=theme.muted)

    highlight_cards = (cover_payload.get("highlight_cards") or [])[:2]
    if not highlight_cards:
        highlight_cards = [
            {"title": "Audience Lens", "body": presentation_json.get("theme_vibe", "Professional & Executive"), "tone": "secondary"},
            {"title": "Deck Style", "body": theme.name, "tone": "primary"},
        ]

    add_card(slide, Inches(6.4), Inches(1.25), Inches(3.6), Inches(2.0), highlight_cards[0].get("title", "Focus"), highlight_cards[0].get("body", ""), theme, accent=accent_for_tone(theme, highlight_cards[0].get("tone")))
    if len(highlight_cards) > 1:
        add_card(slide, Inches(6.4), Inches(3.55), Inches(3.6), Inches(2.0), highlight_cards[1].get("title", "Impact"), highlight_cards[1].get("body", ""), theme, accent=accent_for_tone(theme, highlight_cards[1].get("tone")))


def render_closing_slide(prs: Presentation, presentation_json: Dict[str, Any], theme: ThemeSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    height = prs.slide_height
    apply_background(slide, width, height, theme)

    closing_payload = presentation_json.get("closing_payload") or {}
    add_text_box(slide, Inches(1.0), Inches(2.1), Inches(8.8), Inches(0.4), (closing_payload.get("eyebrow") or "NEXT MOVE").upper(), theme, 12, color=theme.accent, bold=True)
    add_text_box(slide, Inches(1.0), Inches(2.7), Inches(8.8), Inches(1.0), closing_payload.get("headline", "Turn the strategy into execution."), theme, 36, color=theme.text, bold=True, font_name=theme.title_font, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.6), Inches(3.95), Inches(7.6), Inches(0.6), closing_payload.get("subheadline", presentation_json.get("deck_subtitle", "Board-ready narrative and actionable next steps.")), theme, 16, color=theme.muted, align=PP_ALIGN.CENTER)

    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.35), Inches(5.0), Inches(3.3), Inches(0.68))
    pill.fill.solid()
    pill.fill.fore_color.rgb = theme.accent
    pill.line.fill.background()

    tf = pill.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = closing_payload.get("pill_text", presentation_json.get("deck_title", "Executive Strategy Deck"))
    style_paragraph(p, theme, 12, color=theme.inverse_text, bold=True, align=PP_ALIGN.CENTER)


def build_pptx(presentation_json: Dict[str, Any], org_name: str = "Enterprise", theme_vibe: str = "Corporate") -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.author = "OmniPitchAI"
    prs.core_properties.company = "Aisynch Labs"
    theme = resolve_theme(presentation_json.get("theme_vibe") or theme_vibe)

    if not presentation_json.get("deck_title"):
        presentation_json = {
            **presentation_json,
            "deck_title": org_name or "Executive Strategy Deck",
            "deck_subtitle": presentation_json.get("deck_subtitle") or "Strategic narrative and executive summary",
        }

    prs.core_properties.title = presentation_json.get("deck_title", org_name or "Executive Strategy Deck")
    prs.core_properties.subject = presentation_json.get("deck_subtitle", "Executive narrative deck")

    render_cover_slide(prs, presentation_json, theme)

    for index, slide_data in enumerate(presentation_json.get("slides", []), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        render_content_slide(slide, slide_data, theme, prs.slide_width, prs.slide_height, index + 1)

    render_closing_slide(prs, presentation_json, theme)

    fd, path = tempfile.mkstemp(prefix="omnipitch_deck_", suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path
